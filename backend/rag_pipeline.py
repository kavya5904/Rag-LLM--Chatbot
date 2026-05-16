from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
import json
from pathlib import Path
import re
from typing import Any
from urllib.error import URLError
from urllib.request import Request, urlopen
from uuid import uuid4

import chromadb
from chromadb.config import Settings as ChromaSettings
from docx import Document as DocxDocument
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer

try:
    from transformers import pipeline as transformers_pipeline
except Exception:  # pragma: no cover
    transformers_pipeline = None

from config import settings


PROMPT_TEMPLATE = """You are a helpful assistant. Answer only from the context. If not found, say 'Not found in document'.

Context:
{context}
"""


@dataclass
class RetrievedChunk:
    source: str
    excerpt: str
    page: int | None


class RAGService:
    def __init__(self) -> None:
        settings.resolved_vector_store_dir.mkdir(parents=True, exist_ok=True)
        settings.resolved_upload_dir.mkdir(parents=True, exist_ok=True)

        self.rag_mode = settings.rag_mode.lower().strip()
        self.client: OpenAI | None = None
        self.local_embedder: SentenceTransformer | None = None
        self.local_generator: Any | None = None
        self.local_generator_load_failed = False

        if self.rag_mode == "openai":
            if not settings.openai_api_key:
                raise ValueError("OPENAI_API_KEY is required when RAG_MODE=openai")
            client_kwargs: dict[str, Any] = {"api_key": settings.openai_api_key}
            if settings.openai_base_url:
                client_kwargs["base_url"] = settings.openai_base_url
                # Required by OpenRouter (and harmless for other providers)
                client_kwargs["default_headers"] = {
                    "HTTP-Referer": "http://localhost:5173",
                    "X-Title": "PDF Chatbot",
                }
            self.client = OpenAI(**client_kwargs)
        else:
            # Fully local, free embeddings and answer generation.
            self.local_embedder = SentenceTransformer(settings.local_embedding_model)

        self.vector_client = chromadb.PersistentClient(
            path=str(settings.resolved_vector_store_dir),
            settings=ChromaSettings(anonymized_telemetry=False),
        )
        self.collection = self.vector_client.get_or_create_collection(name="pdf_chunks")

    def ingest_document(self, filename: str, content: bytes, user_id: int | None = None) -> int:
        # Store in per-user subdirectory
        if user_id is not None:
            user_upload_dir = settings.resolved_upload_dir / str(user_id)
        else:
            user_upload_dir = settings.resolved_upload_dir
        user_upload_dir.mkdir(parents=True, exist_ok=True)
        output_path = user_upload_dir / filename
        output_path.write_bytes(content)

        text_by_page = self._extract_text_by_unit(filename, content)
        chunks: list[dict[str, Any]] = []
        chunk_index = 0
        owner = str(user_id) if user_id is not None else "__global__"
        for page_number, page_text in text_by_page:
            for chunk in self._chunk_text(page_text):
                chunks.append(
                    {
                        "id": str(uuid4()),
                        "document": chunk,
                        "metadata": {
                            "source": filename,
                            "page": page_number,
                            "chunk_index": chunk_index,
                            "user_id": owner,
                        },
                    }
                )
                chunk_index += 1

        if not chunks:
            raise ValueError("No extractable text found in the document.")

        # Delete existing chunks for this user+file combo
        existing = self.collection.get(where={"$and": [{"source": filename}, {"user_id": owner}]}, include=[])
        existing_ids = existing.get("ids", [])
        if existing_ids:
            self.collection.delete(ids=existing_ids)

        embeddings = self._embed_texts([chunk["document"] for chunk in chunks])
        self.collection.add(
            ids=[chunk["id"] for chunk in chunks],
            documents=[chunk["document"] for chunk in chunks],
            metadatas=[chunk["metadata"] for chunk in chunks],
            embeddings=embeddings,
        )
        return len(chunks)

    def _extract_text_by_unit(self, filename: str, content: bytes) -> list[tuple[int, str]]:
        extension = Path(filename).suffix.lower()
        if extension == ".pdf":
            return self._extract_text_by_page(content)
        if extension == ".docx":
            return self._extract_text_by_docx_section(content)
        if extension in {".pptx", ".ppt"}:
            return self._extract_text_by_slide(content, extension)
        raise ValueError("Unsupported file type. Please upload PDF, DOCX, or PPT/PPTX.")

    def _extract_text_by_docx_section(self, content: bytes) -> list[tuple[int, str]]:
        document = DocxDocument(BytesIO(content))
        paragraphs = [" ".join(paragraph.text.split()) for paragraph in document.paragraphs if paragraph.text.strip()]
        if not paragraphs:
            return []

        # Batch paragraphs into section-sized units for chunking consistency.
        grouped: list[tuple[int, str]] = []
        section_size = 10
        for index in range(0, len(paragraphs), section_size):
            section_no = (index // section_size) + 1
            grouped.append((section_no, "\n".join(paragraphs[index : index + section_size])))
        return grouped

    def _extract_text_by_slide(self, content: bytes, extension: str) -> list[tuple[int, str]]:
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as exc:
            if extension == ".ppt":
                raise ValueError("Legacy .ppt files may fail to parse. Please save as .pptx and re-upload.") from exc
            raise ValueError("Unable to read presentation content.") from exc

        slides_text: list[tuple[int, str]] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    parts.append(" ".join(text.split()))
            merged = "\n".join(parts).strip()
            if merged:
                slides_text.append((slide_index, merged))
        return slides_text

    def answer_question(
        self,
        question: str,
        history: list[dict[str, str]],
        document_name: str | None = None,
        user_id: int | None = None,
    ) -> tuple[str, list[RetrievedChunk]]:
        owner = str(user_id) if user_id is not None else None
        summary_request = self._looks_like_summary_request(question)
        retrieval_k = max(settings.retrieval_k, 8) if summary_request else settings.retrieval_k
        if summary_request:
            relevant_chunks = self.retrieve_for_summary(document_name=document_name, limit=max(retrieval_k, 12), owner=owner)
            if not relevant_chunks:
                relevant_chunks = self.retrieve(question, retrieval_k, document_name=document_name, owner=owner)
        else:
            relevant_chunks = self.retrieve(question, retrieval_k, document_name=document_name, owner=owner)
        if not relevant_chunks:
            return "Not found in document", []

        if summary_request:
            return self._summarize_document(question, history, relevant_chunks), relevant_chunks

        if self.rag_mode == "ollama":
            return self._answer_question_ollama(question, history, relevant_chunks), relevant_chunks

        if self.rag_mode != "openai":
            return self._answer_question_local(question, history, relevant_chunks), relevant_chunks

        context = "\n\n".join(
            f"Source: {chunk.source} | Page: {chunk.page or 'n/a'}\n{chunk.excerpt}" for chunk in relevant_chunks
        )
        messages = [{"role": "system", "content": PROMPT_TEMPLATE.format(context=context or "No matching context found.")}] 
        messages.extend(history[-settings.history_window :])
        messages.append({"role": "user", "content": question})

        response = self.client.chat.completions.create(
            model=settings.openai_chat_model,
            temperature=0.1,
            messages=messages,
        )
        answer = response.choices[0].message.content or "Not found in document"
        return answer.strip(), relevant_chunks

    def _summarize_document(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str:
        if self.rag_mode == "ollama":
            answer = self._summarize_with_ollama(question, history, relevant_chunks)
            if answer:
                return answer

        answer = self._summarize_with_local_generator(question, history, relevant_chunks)
        if answer:
            return answer

        return self._summarize_with_fallback(relevant_chunks)

    def _summarize_with_ollama(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str | None:
        partial_summaries: list[str] = []
        for context in self._build_summary_batches(relevant_chunks):
            prompt = (
                "Summarize the following part of a PDF. Focus on the main ideas, important findings, "
                "and key details. Keep it factual and concise.\n\n"
                f"Context:\n{context}\n\n"
                "Write 3-5 bullet points."
            )
            response = self._call_ollama(prompt)
            if response:
                partial_summaries.append(response)

        if not partial_summaries:
            return None

        recent_history = "\n".join(
            f"{message['role']}: {message['content']}"
            for message in history[-settings.history_window :]
        )
        combine_prompt = (
            "You are summarizing an entire PDF from partial summaries. "
            "Create a clean final answer with these sections exactly: Overview, Key Points, Important Details. "
            "Keep it grounded in the provided summaries.\n\n"
            f"Conversation:\n{recent_history or 'None'}\n\n"
            f"User request: {question}\n\n"
            "Partial summaries:\n"
            + "\n\n".join(partial_summaries)
        )
        return self._call_ollama(combine_prompt)

    def _summarize_with_local_generator(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str | None:
        generator = self._get_local_generator()
        if generator is None:
            return None

        partial_summaries: list[str] = []
        for context in self._build_summary_batches(relevant_chunks):
            prompt = (
                "Summarize this part of a PDF into 3 to 5 concise bullet points. "
                "Keep only the most important information.\n\n"
                f"Context:\n{context}"
            )
            response = self._run_local_generator(prompt, max_new_tokens=160)
            if response:
                partial_summaries.append(response)

        if not partial_summaries:
            return None

        recent_history = "\n".join(
            f"{message['role']}: {message['content']}"
            for message in history[-settings.history_window :]
        )
        combine_prompt = (
            "Create a final PDF summary from these partial summaries. "
            "Use sections named Overview, Key Points, and Important Details. "
            "Be clear and coherent.\n\n"
            f"Conversation:\n{recent_history or 'None'}\n\n"
            f"User request: {question}\n\n"
            "Partial summaries:\n"
            + "\n\n".join(partial_summaries)
        )
        return self._run_local_generator(combine_prompt, max_new_tokens=220)

    def _summarize_with_fallback(self, relevant_chunks: list[RetrievedChunk]) -> str:
        sentences = self._collect_candidate_sentences(relevant_chunks)
        if not sentences:
            return "Not found in document"

        overview = sentences[:2]
        key_points = sentences[2:6]
        if not key_points:
            key_points = sentences[:4]

        lines = ["Overview:"]
        for sentence in overview:
            lines.append(f"- {sentence}")

        lines.append("")
        lines.append("Key Points:")
        for sentence in key_points:
            lines.append(f"- {sentence}")

        return "\n".join(lines).strip()

    def _answer_question_ollama(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str:
        context = self._build_local_context(relevant_chunks)
        if not context:
            return "Not found in document"

        recent_history = "\n".join(
            f"{message['role']}: {message['content']}"
            for message in history[-settings.history_window :]
        )
        prompt = (
            "You are a helpful PDF assistant. Answer only from the provided context. "
            "If the question asks for a summary, provide a clear structured summary with key points. "
            "If the answer is not in the context, say 'Not found in document'.\n\n"
            f"Conversation:\n{recent_history or 'None'}\n\n"
            f"Context:\n{context}\n\n"
            f"Question: {question}\n\n"
            "Answer:"
        )

        try:
            body = json.dumps(
                {
                    "model": settings.ollama_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.2,
                    },
                }
            ).encode("utf-8")
            request = Request(
                url=f"{settings.ollama_base_url.rstrip('/')}/api/generate",
                data=body,
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urlopen(request, timeout=settings.ollama_timeout_seconds) as response:
                payload = json.loads(response.read().decode("utf-8"))
        except (URLError, TimeoutError, OSError, json.JSONDecodeError):
            return self._answer_question_local(question, history, relevant_chunks)

        answer = str(payload.get("response", "")).strip()
        if not answer:
            return self._answer_question_local(question, history, relevant_chunks)
        return answer

    def _answer_question_local(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str:
        generated_answer = self._generate_answer_with_local_model(question, history, relevant_chunks)
        if generated_answer:
            return generated_answer
        return self._answer_question_local_fallback(question, relevant_chunks)

    def _run_local_generator(self, prompt: str, max_new_tokens: int) -> str | None:
        generator = self._get_local_generator()
        if generator is None:
            return None
        try:
            result = generator(
                prompt,
                max_new_tokens=max_new_tokens,
                do_sample=False,
                truncation=True,
            )
        except Exception:
            self.local_generator_load_failed = True
            self.local_generator = None
            return None

        generated_text = (result[0].get("generated_text") or "").strip()
        return generated_text or None

    def _generate_answer_with_local_model(
        self,
        question: str,
        history: list[dict[str, str]],
        relevant_chunks: list[RetrievedChunk],
    ) -> str | None:
        context = self._build_local_context(relevant_chunks)
        if not context:
            return None

        recent_history = "\n".join(
            f"{message['role']}: {message['content']}"
            for message in history[-settings.history_window :]
        )
        prompt = (
            "You are a helpful PDF assistant. Use only the provided context. "
            "If the user asks for a summary, write a concise summary with the main points. "
            "If the answer is not supported by the context, say 'Not found in document'.\n\n"
            f"Conversation:\n{recent_history or 'None'}\n\n"
            f"Context:\n{context}\n\n"
            f"Question: {question}\n\n"
            "Answer:"
        )
        return self._run_local_generator(prompt, max_new_tokens=220)

    def _call_ollama(self, prompt: str) -> str | None:
        try:
            body = json.dumps(
                {
                    "model": settings.ollama_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                        "temperature": 0.2,
                    },
                }
            ).encode("utf-8")
            request = Request(
                url=f"{settings.ollama_base_url.rstrip('/')}/api/generate",
                data=body,
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urlopen(request, timeout=settings.ollama_timeout_seconds) as response:
                payload = json.loads(response.read().decode("utf-8"))
        except (URLError, TimeoutError, OSError, json.JSONDecodeError):
            return None

        answer = str(payload.get("response", "")).strip()
        return answer or None

    def _get_local_generator(self) -> Any | None:
        if self.local_generator_load_failed:
            return None
        if self.local_generator is not None:
            return self.local_generator
        if transformers_pipeline is None:
            self.local_generator_load_failed = True
            return None

        try:
            self.local_generator = transformers_pipeline(
                "text2text-generation",
                model=settings.local_generation_model,
            )
        except Exception:
            self.local_generator_load_failed = True
            self.local_generator = None
            return None
        return self.local_generator

    def _answer_question_local_fallback(self, question: str, relevant_chunks: list[RetrievedChunk]) -> str:
        sentences = self._collect_candidate_sentences(relevant_chunks)
        if not sentences:
            return "Not found in document"

        summary_request = self._looks_like_summary_request(question)
        top_sentences = self._rank_sentences(question, sentences, limit=5 if summary_request else 3)
        if not top_sentences:
            return "Not found in document"

        if summary_request:
            bullet_points = [f"- {sentence}" for sentence in top_sentences]
            return "Summary:\n" + "\n".join(bullet_points)

        return " ".join(top_sentences)

    def _collect_candidate_sentences(self, relevant_chunks: list[RetrievedChunk]) -> list[str]:
        sentences: list[str] = []
        seen: set[str] = set()
        for chunk in relevant_chunks:
            for sentence in re.split(r"(?<=[.!?])\s+", chunk.excerpt):
                cleaned = " ".join(sentence.split()).strip()
                if len(cleaned) < 35:
                    continue
                lowered = cleaned.lower()
                if lowered in seen:
                    continue
                seen.add(lowered)
                sentences.append(cleaned)
        return sentences

    def _build_summary_batches(self, relevant_chunks: list[RetrievedChunk]) -> list[str]:
        batches: list[str] = []
        current_parts: list[str] = []
        current_length = 0
        max_chars = max(1200, settings.local_generation_max_context_chars // 2)

        for chunk in relevant_chunks:
            part = f"Page {chunk.page or 'n/a'}:\n{chunk.excerpt.strip()}"
            if current_parts and current_length + len(part) > max_chars:
                batches.append("\n\n".join(current_parts))
                current_parts = []
                current_length = 0
            current_parts.append(part)
            current_length += len(part)

        if current_parts:
            batches.append("\n\n".join(current_parts))
        return batches

    def _rank_sentences(self, question: str, sentences: list[str], limit: int) -> list[str]:
        question_terms = self._tokenize_for_ranking(question)
        scored: list[tuple[float, str]] = []
        for index, sentence in enumerate(sentences):
            sentence_terms = self._tokenize_for_ranking(sentence)
            overlap = len(question_terms.intersection(sentence_terms))
            density_bonus = min(len(sentence_terms) / 40, 1)
            position_bonus = max(0, 1 - index * 0.03)
            score = overlap * 3 + density_bonus + position_bonus
            scored.append((score, sentence))

        scored.sort(key=lambda item: item[0], reverse=True)
        selected = [sentence for _, sentence in scored[:limit]]
        return selected

    def _tokenize_for_ranking(self, text: str) -> set[str]:
        return {token for token in re.findall(r"[A-Za-z0-9]+", text.lower()) if len(token) > 2}

    def _looks_like_summary_request(self, question: str) -> bool:
        lowered = question.lower()
        return any(
            keyword in lowered
            for keyword in ["summary", "summarize", "overview", "main points", "key points", "brief"]
        )

    def _build_local_context(self, relevant_chunks: list[RetrievedChunk]) -> str:
        sections: list[str] = []
        total_chars = 0
        for chunk in relevant_chunks:
            section = f"Source: {chunk.source} | Page: {chunk.page or 'n/a'}\n{chunk.excerpt.strip()}"
            if total_chars + len(section) > settings.local_generation_max_context_chars:
                remaining = settings.local_generation_max_context_chars - total_chars
                if remaining > 200:
                    sections.append(section[:remaining].rstrip())
                break
            sections.append(section)
            total_chars += len(section)
        return "\n\n".join(sections)

    def retrieve(self, query: str, top_k: int, document_name: str | None = None, owner: str | None = None) -> list[RetrievedChunk]:
        query_kwargs: dict[str, Any] = {
            "query_embeddings": [self._embed_text(query)],
            "n_results": top_k,
            "include": ["documents", "metadatas"],
        }
        where_clauses: list[dict] = []
        if document_name:
            where_clauses.append({"source": document_name})
        if owner:
            where_clauses.append({"user_id": owner})
        if len(where_clauses) > 1:
            query_kwargs["where"] = {"$and": where_clauses}
        elif len(where_clauses) == 1:
            query_kwargs["where"] = where_clauses[0]

        results = self.collection.query(**query_kwargs)
        documents = results.get("documents", [[]])[0]
        metadatas = results.get("metadatas", [[]])[0]

        retrieved: list[RetrievedChunk] = []
        for document, metadata in zip(documents, metadatas, strict=False):
            retrieved.append(
                RetrievedChunk(
                    source=str(metadata.get("source", "Unknown")),
                    page=metadata.get("page"),
                    excerpt=document,
                )
            )
        return retrieved

    def retrieve_for_summary(self, document_name: str | None, limit: int, owner: str | None = None) -> list[RetrievedChunk]:
        target_document = document_name or self._get_single_indexed_document_name(owner=owner)
        if not target_document:
            return []

        where_clauses: list[dict] = [{"source": target_document}]
        if owner:
            where_clauses.append({"user_id": owner})
        where_filter = {"$and": where_clauses} if len(where_clauses) > 1 else where_clauses[0]

        results = self.collection.get(
            where=where_filter,
            include=["documents", "metadatas"],
        )
        documents = results.get("documents", [])
        metadatas = results.get("metadatas", [])

        combined: list[tuple[tuple[int, int], RetrievedChunk]] = []
        for document, metadata in zip(documents, metadatas, strict=False):
            page = metadata.get("page")
            chunk_index = metadata.get("chunk_index", 10**9)
            combined.append(
                (
                    (page or 10**9, chunk_index),
                    RetrievedChunk(
                        source=str(metadata.get("source", "Unknown")),
                        page=page,
                        excerpt=document,
                    ),
                )
            )

        combined.sort(key=lambda item: item[0])
        return [chunk for _, chunk in combined[:limit]]

    def _get_single_indexed_document_name(self, owner: str | None = None) -> str | None:
        get_kwargs: dict[str, Any] = {"include": ["metadatas"]}
        if owner:
            get_kwargs["where"] = {"user_id": owner}
        results = self.collection.get(**get_kwargs)
        metadatas = results.get("metadatas", [])
        document_names = {
            str(metadata.get("source"))
            for metadata in metadatas
            if metadata and metadata.get("source")
        }
        if len(document_names) == 1:
            return next(iter(document_names))
        return None

    def _extract_text_by_page(self, content: bytes) -> list[tuple[int, str]]:
        reader = PdfReader(BytesIO(content))
        pages: list[tuple[int, str]] = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append((index, text))
        return pages

    def _chunk_text(self, text: str) -> list[str]:
        words = text.split()
        if not words:
            return []

        chunk_size = settings.chunk_size
        overlap = settings.chunk_overlap
        chunks: list[str] = []
        start = 0
        while start < len(words):
            end = min(len(words), start + chunk_size)
            chunk = " ".join(words[start:end]).strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(words):
                break
            start = max(0, end - overlap)
        return chunks

    def _embed_text(self, text: str) -> list[float]:
        if self.rag_mode == "openai":
            response = self.client.embeddings.create(model=settings.openai_embedding_model, input=text)
            return response.data[0].embedding
        embedding = self.local_embedder.encode(text, normalize_embeddings=True)
        return embedding.tolist()

    def _embed_texts(self, texts: list[str]) -> list[list[float]]:
        if self.rag_mode == "openai":
            response = self.client.embeddings.create(model=settings.openai_embedding_model, input=texts)
            return [item.embedding for item in response.data]
        embeddings = self.local_embedder.encode(texts, normalize_embeddings=True)
        return embeddings.tolist()


rag_service = RAGService()
